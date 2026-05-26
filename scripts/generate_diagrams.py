"""
generate_diagrams.py
Generates System Architecture, DFD, Flowchart, and Use Case diagrams
Outputs: PNG images (via matplotlib) + fully editable PPTX shapes (via python-pptx).
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = Path(__file__).resolve().parents[1]
OUTPUT_DIR = BASE_DIR / "Project_Diagrams_PNG"
PPTX_PATH = BASE_DIR / "Project_Diagrams.pptx"
OUTPUT_DIR.mkdir(exist_ok=True)

COLORS = {
    "bg": "#FAFBFC", "primary": "#2C3E50", "secondary": "#34495E",
    "accent1": "#2980B9", "accent2": "#27AE60", "accent3": "#E74C3C",
    "accent4": "#F39C12", "accent5": "#8E44AD", "accent6": "#1ABC9C",
    "light_blue": "#D6EAF8", "light_green": "#D5F5E3", "light_red": "#FADBD8",
    "light_orange": "#FDEBD0", "light_purple": "#E8DAEF", "light_teal": "#D1F2EB",
    "light_gray": "#F2F3F4", "border": "#BDC3C7", "white": "#FFFFFF", "line": "#7F8C8D",
}

def hex_to_rgb(hex_c):
    h = hex_c.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def hex_to_rgb_tuple(hex_c):
    h = hex_c.lstrip('#')
    return tuple(int(h[i:i+2], 16)/255.0 for i in (0,2,4))

class MatplotlibCanvas:
    def __init__(self, ax):
        self.ax = ax
    def add_title(self, text, rx, ry, fontsize=15, color=COLORS["primary"], sub=None):
        self.ax.text(rx, ry, text, ha="center", va="center", fontsize=fontsize,
                     fontweight="bold", color=color, family="sans-serif")
        if sub:
            self.ax.text(rx, ry-4, sub, ha="center", va="center", fontsize=10,
                         color=COLORS["secondary"], family="sans-serif", fontstyle="italic")
    def draw_box(self, x, y, w, h, text, fill_color, text_color="#2C3E50",
                 fontsize=9, edge_color="#BDC3C7", style="round,pad=0.02",
                 box_style=None):
        alpha = 0.5 if "dashed" in str(box_style) else 1.0
        ls = "--" if "dashed" in str(box_style) else "-"
        box = FancyBboxPatch((x - w/2, y - h/2), w, h, boxstyle=style,
                             facecolor=fill_color, edgecolor=edge_color,
                             linewidth=1.2, alpha=alpha, linestyle=ls, zorder=2)
        self.ax.add_patch(box)
        if text:
            # Add bbox behind text if needed
            bbox = dict(boxstyle="round,pad=0.3", fc="white", ec=edge_color, lw=1) if "paddedText" in str(box_style) else None
            self.ax.text(x, y, text, ha="center", va="center", fontsize=fontsize,
                         fontweight="bold", color=text_color, zorder=3,
                         wrap=True, family="sans-serif", linespacing=1.3, bbox=bbox)
    def draw_circle(self, x, y, radius, text, fill_color, text_color="#FFFFFF",
                    fontsize=8, edge_color="#BDC3C7", hw=1):
        circle = plt.Circle((x, y), radius, facecolor=fill_color,
                            edgecolor=edge_color, linewidth=1.2, zorder=2)
        self.ax.add_patch(circle)
        self.ax.text(x, y, text, ha="center", va="center", fontsize=fontsize,
                     fontweight="bold", color=text_color, zorder=3, family="sans-serif")
    def draw_arrow(self, x1, y1, x2, y2, color="#7F8C8D", label="",
                   style="->", label_offset=(0,4), fontsize=7):
        self.ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                         arrowprops=dict(arrowstyle=style, color=color, lw=1.5), zorder=1)
        if label:
            mx, my = (x1+x2)/2 + label_offset[0], (y1+y2)/2 + label_offset[1]
            self.ax.text(mx, my, label, fontsize=fontsize, ha="center", va="center",
                         color=COLORS["secondary"], family="sans-serif", fontstyle="italic",
                         bbox=dict(boxstyle="round,pad=0.15", fc="white", ec="none", alpha=0.85), zorder=4)
    def draw_diamond(self, x, y, w, h, text, fill_color, text_color="#2C3E50", fontsize=7.5, edge_color="#BDC3C7"):
        diamond = plt.Polygon([(x, y+h/2), (x+w/2, y), (x, y-h/2), (x-w/2, y)],
                              facecolor=fill_color, edgecolor=edge_color, linewidth=1.2, zorder=2)
        self.ax.add_patch(diamond)
        self.ax.text(x, y, text, ha="center", va="center", fontsize=fontsize, fontweight="bold",
                     color=text_color, zorder=3, family="sans-serif")
    def draw_parallelogram(self, x, y, w, h, text, fill_color, text_color="#2C3E50", fontsize=8, edge_color="#BDC3C7"):
        skew = w * 0.12
        para = plt.Polygon([(x-w/2+skew, y+h/2), (x+w/2+skew, y+h/2),
                            (x+w/2-skew, y-h/2), (x-w/2-skew, y-h/2)],
                           facecolor=fill_color, edgecolor=edge_color, linewidth=1.2, zorder=2)
        self.ax.add_patch(para)
        self.ax.text(x, y, text, ha="center", va="center", fontsize=fontsize, fontweight="bold",
                     color=text_color, zorder=3, family="sans-serif")
    def draw_rounded_stadium(self, x, y, w, h, text, fill_color, text_color="#FFFFFF", fontsize=8, edge_color="#BDC3C7"):
        box = FancyBboxPatch((x-w/2, y-h/2), w, h, boxstyle=f"round,pad={h/2}", facecolor=fill_color, edgecolor=edge_color, linewidth=1.5, zorder=2)
        self.ax.add_patch(box)
        self.ax.text(x, y, text, ha="center", va="center", fontsize=fontsize, fontweight="bold", color=text_color, zorder=3, family="sans-serif")
    def draw_actor(self, x, y, label, fontsize=8):
        self.ax.add_patch(plt.Circle((x, y+16), 3.5, facecolor=COLORS["white"], edgecolor=COLORS["primary"], lw=1.5, zorder=3))
        self.ax.plot([x, x], [y+12.5, y+3], color=COLORS["primary"], lw=1.5, zorder=3)
        self.ax.plot([x-7, x+7], [y+9, y+9], color=COLORS["primary"], lw=1.5, zorder=3)
        self.ax.plot([x, x-6], [y+3, y-5], color=COLORS["primary"], lw=1.5, zorder=3)
        self.ax.plot([x, x+6], [y+3, y-5], color=COLORS["primary"], lw=1.5, zorder=3)
        self.ax.text(x, y-12, label, ha="center", va="top", fontsize=fontsize, fontweight="bold", color=COLORS["primary"], family="sans-serif")

class PptxCanvas:
    def __init__(self, slide, center_x, center_y, scale=0.063, sw=13.333, sh=7.5):
        self.slide = slide
        self.cx = center_x
        self.cy = center_y
        self.scale = scale
        self.sw = sw
        self.sh = sh
    def _tr(self, x, y, w=0, h=0):
        px = (x - self.cx) * self.scale + (self.sw/2)
        py = -(y - self.cy) * self.scale + (self.sh/2)
        pw, ph = w * self.scale, h * self.scale
        return Inches(px - pw/2), Inches(py - ph/2), Inches(pw), Inches(ph)
    def _cPoint(self, x, y):
        px = (x - self.cx) * self.scale + (self.sw/2)
        py = -(y - self.cy) * self.scale + (self.sh/2)
        return Inches(px), Inches(py)
    def _fmt(self, shape, t, f, tc, ec, fs, is_dashed=False):
        if f:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(f)
            if is_dashed: shape.fill.transparency = 0.5
        else:
            shape.fill.background()
        if ec and ec != "none":
            shape.line.color.rgb = hex_to_rgb(ec)
            shape.line.width = Pt(1.5)
            from pptx.enum.dml import MSO_LINE
            if is_dashed: shape.line.dash_style = MSO_LINE.DASH
        else:
            shape.line.fill.background()
        if t:
            tf = shape.text_frame
            tf.word_wrap = True
            try:
                from pptx.enum.text import MSO_ANCHOR
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except: pass
            tf.clear()
            p = tf.paragraphs[0]
            p.text = t
            p.font.color.rgb = hex_to_rgb(tc)
            p.font.size = Pt(fs)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    def add_title(self, text, rx, ry, fontsize=15, color=COLORS["primary"], sub=None):
        pass # Titles rendered in slide layouts instead
    def draw_box(self, x, y, w, h, text, fill_color, text_color="#2C3E50", fontsize=9, edge_color="#BDC3C7", style="", box_style=None):
        l, t, wd, ht = self._tr(x,y,w,h)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, wd, ht)
        is_dashed = "dashed" in str(box_style)
        self._fmt(shape, text, fill_color, text_color, edge_color, fontsize, is_dashed)
        # Note: PPTX text box over dashed rectangle... PPTX keeps text centered in shape
    def draw_circle(self, x, y, radius, text, fill_color, text_color="#FFFFFF", fontsize=8, edge_color="#BDC3C7", hw=1):
        l, t, wd, ht = self._tr(x,y,radius*hw*2,radius*2)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, wd, ht)
        self._fmt(shape, text, fill_color, text_color, edge_color, fontsize)
    def draw_diamond(self, x, y, w, h, text, fill_color, text_color="#2C3E50", fontsize=7.5, edge_color="#BDC3C7"):
        l, t, wd, ht = self._tr(x,y,w,h)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, wd, ht)
        self._fmt(shape, text, fill_color, text_color, edge_color, fontsize)
    def draw_parallelogram(self, x, y, w, h, text, fill_color, text_color="#2C3E50", fontsize=8, edge_color="#BDC3C7"):
        l, t, wd, ht = self._tr(x,y,w,h)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, l, t, wd, ht)
        self._fmt(shape, text, fill_color, text_color, edge_color, fontsize)
    def draw_rounded_stadium(self, x, y, w, h, text, fill_color, text_color="#FFFFFF", fontsize=8, edge_color="#BDC3C7"):
        l, t, wd, ht = self._tr(x,y,w,h)
        shape = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, wd, ht)
        self._fmt(shape, text, fill_color, text_color, edge_color, fontsize)
    def draw_arrow(self, x1, y1, x2, y2, color="#7F8C8D", label="", style="->", label_offset=(0,4), fontsize=7):
        bx, by = self._cPoint(x1, y1)
        ex, ey = self._cPoint(x2, y2)
        conn = self.slide.shapes.add_connector(1, bx, by, ex, ey)
        conn.line.color.rgb = hex_to_rgb(color)
        conn.line.width = Pt(1.5)
        if "->" in style:
            conn.line.end_arrowhead = 2
        
        if label:
            mx, my = self._cPoint((x1+x2)/2, (y1+y2)/2)
            dx = Inches(label_offset[0]*self.scale)
            dy = Inches(label_offset[1]*self.scale)
            sh = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, mx+dx-Inches(0.6), my-dy-Inches(0.2), Inches(1.2), Inches(0.4))
            self._fmt(sh, label, "#FFFFFF", "#34495E", "none", fontsize)
    def draw_actor(self, x, y, label, fontsize=8):
        head_r = 3.5
        self.draw_circle(x, y+16, head_r, "", None, edge_color=COLORS["primary"])
        self.draw_arrow(x, y+12.5, x, y+3, COLORS["primary"], style="-")
        self.draw_arrow(x-7, y+9, x+7, y+9, COLORS["primary"], style="-")
        self.draw_arrow(x, y+3, x-6, y-5, COLORS["primary"], style="-")
        self.draw_arrow(x, y+3, x+6, y-5, COLORS["primary"], style="-")
        l, t = self._cPoint(x, y-12)
        sh = self.slide.shapes.add_textbox(l-Inches(0.5), t, Inches(1), Inches(0.4))
        self._fmt(sh, label, None, COLORS["primary"], "none", 8)

# DRAW DEFINITIONS
def d_arch(c):
    c.add_title("V5 Hybrid System Architecture", 90, 140, 16)
    
    # Layer 1: Client
    c.draw_box(90, 122, 50, 10, "Streamlit Dashboard\n(User Interface & CSV Upload)", COLORS["light_green"], edge_color=COLORS["accent2"])
    
    # Layer 2: API
    c.draw_arrow(90, 117, 90, 109, COLORS["accent2"], label="API Request")
    c.draw_box(90, 104, 50, 10, "FastAPI Backend\n(/predict endpoint)", COLORS["light_purple"], edge_color=COLORS["accent5"])
    
    # Layer 3: Preprocessing
    c.draw_arrow(90, 99, 90, 91, COLORS["accent5"])
    c.draw_box(90, 86, 50, 10, "Data Preprocessing\n(PaySim Feature Engineering)", COLORS["light_blue"], edge_color=COLORS["accent1"])
    
    # Layer 4: V5 Engine Container
    c.draw_arrow(90, 81, 90, 73, COLORS["accent1"])
    c.draw_box(90, 52, 120, 32, "", "#F8F9FA", edge_color=COLORS["accent1"], box_style="dashed")
    c.draw_box(90, 68, 40, 6, "V5 Hybrid Model Engine", "#FFFFFF", COLORS["accent1"], edge_color=COLORS["accent1"])
    
    # Engine Models
    c.draw_arrow(90, 65, 60, 58, COLORS["accent1"])
    c.draw_arrow(90, 65, 120, 58, COLORS["accent1"])
    
    c.draw_box(60, 48, 44, 12, "PATH A (Tier 1)\nXGBoost + Random Forest\n(Known Fraud)", COLORS["light_green"], edge_color=COLORS["accent2"])
    c.draw_box(120, 48, 44, 12, "PATH B (Tier 2/3)\nBiLSTM + Autoencoder\n(Novel Fraud)", COLORS["light_purple"], edge_color=COLORS["accent5"])
    
    # Layer 5: Decision
    c.draw_arrow(60, 42, 90, 31, COLORS["accent2"])
    c.draw_arrow(120, 42, 90, 31, COLORS["accent5"])
    
    c.draw_box(90, 26, 50, 10, "Decision Engine", COLORS["light_gray"], edge_color=COLORS["primary"])
    
    # Layer 6: Outputs
    c.draw_arrow(90, 21, 30, 10, COLORS["accent3"], label="Tier 1")
    c.draw_arrow(90, 21, 70, 10, COLORS["accent5"], label="Tier 2")
    c.draw_arrow(90, 21, 110, 10, COLORS["accent4"], label="Tier 3")
    c.draw_arrow(90, 21, 150, 10, COLORS["accent2"], label="Tier 4")
    
    c.draw_box(30, 5, 26, 10, "BLOCK\n(Known Fraud)", COLORS["light_red"], edge_color=COLORS["accent3"])
    c.draw_box(70, 5, 26, 10, "BLOCK_NOVEL\n(Novel Fraud)", COLORS["light_purple"], edge_color=COLORS["accent5"])
    c.draw_box(110, 5, 26, 10, "REVIEW\n(Suspicious)", COLORS["light_orange"], edge_color=COLORS["accent4"])
    c.draw_box(150, 5, 26, 10, "ALLOW\n(Legitimate)", COLORS["light_green"], edge_color=COLORS["accent2"])

def d_dfd0(c):
    c.add_title("DFD Level 0 — Context Diagram", 80, 101, 15, sub="Fraud Detection System Boundary")
    c.draw_box(10, 70, 24, 14, "User /\nFinancial\nInstitution", COLORS["light_blue"], edge_color=COLORS["accent1"])
    c.draw_box(150, 70, 24, 14, "Admin /\nFraud Analyst", COLORS["light_green"], edge_color=COLORS["accent2"])
    c.draw_box(10, 25, 24, 14, "Transaction\nDatabase\n(PaySim)", COLORS["light_orange"], edge_color=COLORS["accent4"])
    c.draw_box(150, 25, 24, 14, "Model\nRepository\n(Trained Models)", COLORS["light_purple"], edge_color=COLORS["accent5"])

    c.draw_circle(80, 48, 22, "0\n\nFraud Detection\nSystem\n(V5 Hybrid)", COLORS["accent1"], edge_color=COLORS["primary"])

    c.draw_arrow(22, 70, 59, 55, COLORS["accent1"], label="Transaction Request")
    c.draw_arrow(59, 42, 22, 30, COLORS["accent1"], label="Processed Data")
    c.draw_arrow(101, 55, 138, 70, COLORS["accent2"], label="Detection Result\n(BLOCK/REVIEW/ALLOW)")
    c.draw_arrow(138, 62, 101, 45, COLORS["accent2"], label="Review Decision")
    c.draw_arrow(22, 32, 59, 42, COLORS["accent4"], label="Historical\nTransaction Data")
    c.draw_arrow(138, 30, 101, 40, COLORS["accent5"], label="Model Artifacts\n(XGB, RF, BiLSTM, AE)")

def d_dfd1(c):
    c.add_title("DFD Level 1 — Fraud Detection System Decomposition", 90, 140, 15)
    c.draw_box(15, 115, 24, 12, "User /\nTransaction\nSource", COLORS["light_blue"], edge_color=COLORS["accent1"])
    c.draw_box(165, 115, 24, 12, "Fraud\nAnalyst /\nAdmin", COLORS["light_green"], edge_color=COLORS["accent2"])
    c.draw_box(15, 25, 24, 12, "D1\nTransaction\nDatasets", COLORS["light_orange"], edge_color=COLORS["accent4"])
    c.draw_box(165, 25, 24, 12, "D2\nModel\nArtifacts", COLORS["light_purple"], edge_color=COLORS["accent5"])

    c.draw_circle(55, 100, 14, "1.0\n\nData\nIngestion &\nValidation", COLORS["accent1"])
    c.draw_circle(125, 100, 14, "2.0\n\nFeature\nEngineering &\nPreprocessing", COLORS["accent2"])
    c.draw_circle(55, 55, 14, "3.0\n\nV5 Hybrid\nModel\nInference", COLORS["accent5"])
    c.draw_circle(125, 55, 14, "4.0\n\nThree-Tier\nDecision\nEngine", COLORS["accent3"])
    c.draw_circle(90, 10, 14, "5.0\n\nResult\nPresentation &\nReporting", COLORS["accent4"])

    c.draw_arrow(27, 115, 42, 107, COLORS["accent1"], label="Raw Transaction")
    c.draw_arrow(69, 100, 111, 100, COLORS["accent1"], label="Validated Data")
    c.draw_arrow(125, 86, 90, 65, COLORS["accent2"], label="Engineered Features")
    c.draw_arrow(90, 62, 55, 65, COLORS["accent2"])
    c.draw_arrow(69, 52, 111, 52, COLORS["accent5"], label="Model Scores &\nAnomaly Flags")
    c.draw_arrow(125, 41, 101, 17, COLORS["accent3"], label="Decision Output")
    # Using straight line for PPTX compatibility
    c.draw_arrow(99, 22, 153, 115, COLORS["accent2"], label="Alerts & Reports", label_offset=(-10, 5))
    c.draw_arrow(27, 25, 45, 42, COLORS["accent4"], label="Historical Data")
    c.draw_arrow(153, 30, 130, 42, COLORS["accent5"], label="Trained Models")

def d_dfd2(c):
    c.add_title("DFD Level 2 — V5 Hybrid Model Inference (Process 3.0 Decomposition)", 100, 151, 14)
    c.draw_box(15, 120, 24, 12, "Engineered\nFeature\nVector", COLORS["light_blue"], edge_color=COLORS["accent1"])
    c.draw_circle(60, 120, 12, "3.1\n\nFeature\nScaling", COLORS["accent1"])

    c.draw_box(57.5, 87.5, 85, 35, "", "#E8F8F5", edge_color=COLORS["accent2"], box_style="dashed")
    c.draw_box(57.5, 102, 35, 6, "PATH A — Known Fraud Detection (V3)", "#E8F8F5", COLORS["accent2"], edge_color="none")
    c.draw_circle(35, 85, 10, "3.2\n\nAE Recon.\nError Calc.", COLORS["accent6"])
    c.draw_circle(75, 85, 10, "3.3\n\nXGB + RF\nEnsemble\n(19 Features)", COLORS["accent2"])

    c.draw_box(153, 87.5, 90, 35, "", "#F5EEF8", edge_color=COLORS["accent5"], box_style="dashed")
    c.draw_box(153, 102, 35, 6, "PATH B — Novel Fraud Detection (V4)", "#F5EEF8", COLORS["accent5"], edge_color="none")
    c.draw_circle(130, 85, 10, "3.4\n\nBiLSTM +\nAttention\nSequence", COLORS["accent5"])
    c.draw_circle(175, 85, 10, "3.5\n\nAnomaly\nDetection\n(AE + IF)", COLORS["accent3"])

    c.draw_circle(100, 42, 14, "3.6\n\nThree-Tier\nClassification", COLORS["accent3"])

    y_tier = 8
    c.draw_box(25, y_tier, 26, 10, "Tier 1\nBLOCK\n(Known Fraud)", COLORS["light_red"], edge_color=COLORS["accent3"])
    c.draw_box(75, y_tier, 26, 10, "Tier 2\nBLOCK_NOVEL\n(Novel Fraud)", COLORS["light_purple"], edge_color=COLORS["accent5"])
    c.draw_box(125, y_tier, 26, 10, "Tier 3\nREVIEW\n(Suspicious)", COLORS["light_orange"], edge_color=COLORS["accent4"])
    c.draw_box(175, y_tier, 26, 10, "Tier 4\nALLOW\n(Legitimate)", COLORS["light_green"], edge_color=COLORS["accent2"])

    c.draw_box(15, 55, 22, 10, "D2\nV3 Model\nArtifacts", COLORS["light_green"], edge_color=COLORS["accent2"])
    c.draw_box(185, 55, 22, 10, "D3\nV4 Model\nArtifacts", COLORS["light_purple"], edge_color=COLORS["accent5"])

    c.draw_arrow(27, 120, 48, 120, COLORS["accent1"], label="Raw Features")
    c.draw_arrow(60, 108, 40, 95, COLORS["accent1"], label="Scaled (V3)")
    c.draw_arrow(60, 108, 125, 95, COLORS["accent1"], label="Scaled (V4)")
    c.draw_arrow(60, 108, 170, 95, COLORS["accent1"])

    c.draw_arrow(45, 85, 65, 85, COLORS["accent6"], label="AE Error")
    c.draw_arrow(75, 75, 90, 53, COLORS["accent2"], label="V3 Score")
    c.draw_arrow(130, 75, 108, 53, COLORS["accent5"], label="Seq Score")
    c.draw_arrow(175, 75, 112, 53, COLORS["accent3"], label="Anomaly Flags")

    c.draw_arrow(26, 55, 30, 76, COLORS["accent2"], label="V3 Weights")
    c.draw_arrow(174, 55, 170, 76, COLORS["accent5"], label="V4 Weights")

    c.draw_arrow(87, 35, 38, 15, COLORS["accent3"])
    c.draw_arrow(93, 30, 75, 15, COLORS["accent5"])
    c.draw_arrow(107, 30, 125, 15, COLORS["accent4"])
    c.draw_arrow(113, 35, 162, 15, COLORS["accent2"])

def d_flow(c):
    c.add_title("Fraud Detection System — Process Flowchart", 80, 222, 15)
    y = 210
    c.draw_rounded_stadium(80, y, 35, 8, "START", COLORS["accent1"], edge_color=COLORS["primary"])
    c.draw_arrow(80, 206, 80, y-11, COLORS["line"])
    y -= 16
    c.draw_parallelogram(80, y, 42, 10, "Receive Transaction\nData (CSV / API)", COLORS["light_blue"], edge_color=COLORS["accent1"])
    c.draw_arrow(80, y-5, 80, y-13, COLORS["line"])
    y -= 18
    # Skip CreditCard check completely
    c.draw_box(80, y, 38, 10, "PaySim Preprocessor\n(Clean, Engineer UPI\n& Velocity Features)", COLORS["light_blue"], edge_color=COLORS["accent1"])
    c.draw_arrow(80, y-5, 80, y-13, COLORS["line"])
    y -= 18
    c.draw_box(80, y, 40, 10, "Standard Scaler\n(Feature Normalization)", COLORS["light_teal"], edge_color=COLORS["accent6"])
    c.draw_arrow(80, y-5, 80, y-13, COLORS["line"])
    y -= 18
    c.draw_diamond(80, y, 48, 14, "Apply V5 Hybrid\nModel?", COLORS["light_gray"], edge_color=COLORS["primary"])

    yp = y - 24
    c.draw_arrow(56, y, 35, yp+5, COLORS["accent2"], label="Path A")
    c.draw_box(35, yp, 36, 10, "V3 Engine\nAutoencoder → AE Error\nXGBoost + RF (19 feat.)", COLORS["light_green"], edge_color=COLORS["accent2"])

    c.draw_arrow(104, y, 125, yp+5, COLORS["accent5"], label="Path B")
    c.draw_box(125, yp, 36, 10, "V4 Engine\nBiLSTM + Attention\nAnomaly (AE + IForest)", COLORS["light_purple"], edge_color=COLORS["accent5"])

    y -= 42
    c.draw_arrow(35, y+13, 65, y+5, COLORS["accent2"])
    c.draw_arrow(125, y+13, 95, y+5, COLORS["accent5"])

    c.draw_diamond(80, y, 48, 14, "V3 Confidence ≥\nBlock Threshold?", COLORS["light_red"], edge_color=COLORS["accent3"])
    c.draw_arrow(56, y-2, 20, y-15, COLORS["accent3"], label="Yes")
    c.draw_box(20, y-22, 28, 9, "TIER 1:\nBLOCK\n(Known Fraud)", COLORS["light_red"], edge_color=COLORS["accent3"])

    c.draw_arrow(80, y-7, 80, y-13, COLORS["line"], label="No")
    y -= 32
    c.draw_diamond(80, y, 52, 14, "SeqScore ≥ Block\nThreshold & Anomaly?", COLORS["light_purple"], edge_color=COLORS["accent5"])
    c.draw_arrow(56, y-2, 20, y-15, COLORS["accent5"], label="Yes")
    c.draw_box(20, y-22, 28, 9, "TIER 2:\nBLOCK_NOVEL\n(Novel Fraud)", COLORS["light_purple"], edge_color=COLORS["accent5"])

    c.draw_arrow(80, y-7, 80, y-13, COLORS["line"], label="No")
    y -= 32
    c.draw_diamond(80, y, 52, 14, "Anomaly Flag OR\nSeqScore ≥ Review?", COLORS["light_orange"], edge_color=COLORS["accent4"])
    c.draw_arrow(56, y-2, 20, y-15, COLORS["accent4"], label="Yes")
    c.draw_box(20, y-22, 28, 9, "TIER 3:\nREVIEW\n(Manual Check)", COLORS["light_orange"], edge_color=COLORS["accent4"])

    c.draw_arrow(80, y-7, 80, y-13, COLORS["line"], label="No")
    y -= 32
    c.draw_box(80, y, 35, 9, "ALLOW\n(Legitimate Transaction)", COLORS["light_green"], edge_color=COLORS["accent2"])
    c.draw_arrow(80, y-4.5, 80, y-11, COLORS["line"])
    y -= 16
    c.draw_parallelogram(80, y, 42, 10, "Return Response\n(Decision + Explanation)", COLORS["light_blue"], edge_color=COLORS["accent1"])
    c.draw_arrow(80, y-5, 80, y-11, COLORS["line"])
    y -= 16
    c.draw_rounded_stadium(80, y, 35, 8, "END", COLORS["accent3"], edge_color=COLORS["primary"])

def d_usecase(c):
    c.add_title("Enhanced Use Case Diagram", 90, 140, 15)
    
    # System Boundary
    c.draw_box(90, 65, 120, 130, "", "#F8F9FA", edge_color=COLORS["primary"], box_style="dashed")
    c.draw_box(90, 130, 48, 8, "V5 Hybrid Fraud Detection System", "#FFFFFF", COLORS["primary"], edge_color="none", fontsize=10)

    # Actors
    c.draw_actor(10, 105, "Transaction\nSource")
    c.draw_actor(10, 35, "System\nAdmin")
    c.draw_actor(170, 105, "Fraud\nAnalyst")
    c.draw_actor(170, 35, "Data\nScientist")

    # Use Cases
    # Use cases for Tx Source
    cx1=65
    c.draw_circle(cx1, 115, 6, "Initiate\nTransaction", COLORS["light_blue"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent1"], hw=3)
    c.draw_circle(cx1, 95, 6, "Upload\nBatch CSV", COLORS["light_blue"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent1"], hw=3)
    
    # Use cases for Admin
    c.draw_circle(cx1, 45, 6, "Manage\nUsers", COLORS["light_purple"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent5"], hw=3)
    c.draw_circle(cx1, 25, 6, "Configure\nThresholds", COLORS["light_purple"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent5"], hw=3)
    
    # Use cases for Fraud Analyst
    cx2=115
    c.draw_circle(cx2, 115, 6, "Review\nTier 3 Alerts", COLORS["light_orange"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent4"], hw=3)
    c.draw_circle(cx2, 95, 6, "Generate\nAudit Reports", COLORS["light_orange"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent4"], hw=3)
    c.draw_circle(cx2, 75, 6, "Analyze Fraud\nTrends", COLORS["light_orange"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent4"], hw=3)
    
    # Use cases for Data Scientist
    c.draw_circle(cx2, 45, 6, "Retrain\nV5 Model", COLORS["light_teal"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent6"], hw=3)
    c.draw_circle(cx2, 25, 6, "Evaluate\nModel Drift", COLORS["light_teal"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent6"], hw=3)

    # Include/Extend Center Nodes
    cx3 = 90
    c.draw_circle(cx3, 105, 6, "Trigger V5\nDetection", COLORS["light_red"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent3"], hw=4)
    c.draw_circle(cx3, 65, 6, "Extract\nFeatures", COLORS["light_green"], text_color=COLORS["primary"], fontsize=6.5, edge_color=COLORS["accent2"], hw=4)

    # Connections
    # Actor to Use Case
    c.draw_arrow(20, 110, cx1-24, 115, COLORS["line"], style="-")
    c.draw_arrow(20, 110, cx1-24, 95, COLORS["line"], style="-")
    
    c.draw_arrow(20, 40, cx1-24, 45, COLORS["line"], style="-")
    c.draw_arrow(20, 40, cx1-24, 25, COLORS["line"], style="-")
    
    c.draw_arrow(160, 110, cx2+24, 115, COLORS["line"], style="-")
    c.draw_arrow(160, 110, cx2+24, 95, COLORS["line"], style="-")
    c.draw_arrow(160, 110, cx2+24, 75, COLORS["line"], style="-")
    
    c.draw_arrow(160, 40, cx2+24, 45, COLORS["line"], style="-")
    c.draw_arrow(160, 40, cx2+24, 25, COLORS["line"], style="-")
    
    # Includes / Extends
    c.draw_arrow(cx1+20, 115, cx3-18, 105, COLORS["accent5"], label="<<include>>", fontsize=5)
    c.draw_arrow(cx1+20, 95, cx3-18, 105, COLORS["accent5"], label="<<include>>", fontsize=5)
    
    c.draw_arrow(cx3+18, 105, cx2-20, 115, COLORS["accent4"], label="<<extend>>", fontsize=5)
    
    c.draw_arrow(cx2-20, 45, cx3+18, 65, COLORS["accent5"], label="<<include>>", fontsize=5)

def d_uml(c):
    c.add_title("UML Class Diagram — V5 Hybrid System", 90, 140, 15)
    
    # Layer 1: Interface / Presentation
    c.draw_box(45, 110, 58, 28, "FraudDashboard (UI)\n" + "—"*22 + "\n+ upload_csv(file)\n+ request_prediction(data)\n+ display_results(json)", COLORS["light_blue"], edge_color=COLORS["accent1"], fontsize=7)
    
    c.draw_box(135, 110, 58, 28, "PredictionAPI (FastAPI)\n" + "—"*22 + "\n+ predict_fraud(json)\n+ return_decision(dict)", COLORS["light_purple"], edge_color=COLORS["accent5"], fontsize=7)
    
    # Layer 2: Business Logic / Preprocessing
    c.draw_box(135, 65, 58, 28, "DataPreprocessor\n" + "—"*22 + "\n+ clean_data(df)\n+ extract_features(df)\n+ scale_features(df)", COLORS["light_green"], edge_color=COLORS["accent2"], fontsize=7)
    
    c.draw_box(45, 65, 58, 28, "V5HybridEngine\n" + "—"*22 + "\n+ detect_fraud(features)\n+ ensemble_predict()", COLORS["light_orange"], edge_color=COLORS["accent4"], fontsize=7)
    
    # Layer 3: Models (Path A & Path B)
    c.draw_box(45, 20, 58, 28, "V3Model (Tier 1)\n" + "—"*22 + "\n+ predict_known()\n- run_xgboost()\n- run_rf()", COLORS["light_red"], edge_color=COLORS["accent3"], fontsize=7)
    
    c.draw_box(135, 20, 58, 28, "V4Model (Tier 2/3)\n" + "—"*22 + "\n+ predict_novel()\n- run_bilstm()\n- get_ae_anomaly()", COLORS["light_teal"], edge_color=COLORS["accent6"], fontsize=7)
    
    # Relationships
    c.draw_arrow(74, 110, 106, 110, COLORS["line"], label="Calls API")
    c.draw_arrow(135, 96, 135, 79, COLORS["line"], label="Uses")
    c.draw_arrow(106, 65, 74, 65, COLORS["line"], label="Feeds Data")
    c.draw_arrow(45, 51, 45, 34, COLORS["line"], label="Delegates")
    
    # Engine to Path B
    c.draw_arrow(74, 55, 106, 34, COLORS["line"], label="Delegates")

diagrams = [
    { "id": 0, "name": "01_System_Architecture", "fn": d_arch, "lim": [(-10,190), (-10,140)], "cx": 90, "cy": 65, "sc": 0.05 },
    { "id": 1, "name": "02_DFD_Level_0", "fn": d_dfd0, "lim": [(-10,170), (-5,105)], "cx": 80, "cy": 50, "sc": 0.065 },
    { "id": 2, "name": "03_DFD_Level_1", "fn": d_dfd1, "lim": [(-10,190), (-10,145)], "cx": 90, "cy": 67.5, "sc": 0.05 },
    { "id": 3, "name": "04_DFD_Level_2", "fn": d_dfd2, "lim": [(-10,210), (-15,155)], "cx": 100, "cy": 70, "sc": 0.045 },
    { "id": 4, "name": "05_Flowchart", "fn": d_flow, "lim": [(-10,170), (-5,225)], "cx": 80, "cy": 110, "sc": 0.035 },
    { "id": 5, "name": "06_Use_Case", "fn": d_usecase, "lim": [(-15,195), (-15,140)], "cx": 90, "cy": 62.5, "sc": 0.045 },
    { "id": 6, "name": "07_UML_Class_Diagram", "fn": d_uml, "lim": [(-10,190), (-10,145)], "cx": 90, "cy": 65, "sc": 0.05 }
]

def main():
    print("="*50)
    print("GENERATING NATIVE DIAGRAMS")
    print("="*50)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    prs_flow = Presentation()
    prs_flow.slide_width = Inches(8.5)
    prs_flow.slide_height = Inches(14.0)
    
    slide_titles = [
        "System Architecture — V5 Hybrid Fraud Detection",
        "DFD Level 0 — Context Diagram",
        "DFD Level 1 — System Decomposition",
        "DFD Level 2 — Model Inference Decomposition",
        "Process Flowchart — Fraud Detection Pipeline",
        "Use Case Diagram — System Interactions",
        "UML Class Diagram — System Components"
    ]

    for d in diagrams:
        print(f"Drawing {d['name']}...")
        # 1. Matplotlib
        fig, ax = plt.subplots(1, 1, figsize=(18, 14))
        fig.patch.set_facecolor(COLORS["bg"])
        ax.set_xlim(d["lim"][0])
        ax.set_ylim(d["lim"][1])
        ax.set_aspect("equal")
        ax.axis("off")
        c_plt = MatplotlibCanvas(ax)
        d["fn"](c_plt)
        fig.savefig(str(OUTPUT_DIR / f"{d['name']}.png"), dpi=200, bbox_inches="tight", facecolor=COLORS["bg"])
        plt.close(fig)

        # 2. Native PPTX
        if d["id"] == 4:
            target_prs = prs_flow
            sw, sh = 8.5, 14.0
        else:
            target_prs = prs
            sw, sh = 13.333, 7.5

        slide = target_prs.slides.add_slide(target_prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
        
        # Add PPTX headers/footers
        ts = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(sw), Inches(0.75))
        ts.fill.solid()
        ts.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        p = ts.text_frame.paragraphs[0]
        p.text = slide_titles[d["id"]]
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.size = Pt(24)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        fs = slide.shapes.add_textbox(Inches(0), Inches(sh - 0.4), Inches(sw), Inches(0.35))
        fp = fs.text_frame.paragraphs[0]
        fp.text = "V5 Hybrid Fraud Detection System  ·  Major Project Documentation"
        fp.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
        fp.font.size = Pt(9)
        fp.alignment = PP_ALIGN.CENTER

        c_pptx = PptxCanvas(slide, d["cx"], d["cy"], d["sc"], sw=sw, sh=sh)
        d["fn"](c_pptx)
        
    prs.save(str(PPTX_PATH))
    flow_path = BASE_DIR / "Project_Flowchart_Word_Format.pptx"
    prs_flow.save(str(flow_path))
    print(f"Done! Main PPTX saved to {PPTX_PATH}")
    print(f"Flowchart PPTX saved to {flow_path}")

if __name__ == "__main__":
    main()
