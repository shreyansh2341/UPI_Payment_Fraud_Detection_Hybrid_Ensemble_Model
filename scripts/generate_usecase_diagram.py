"""
generate_usecase_diagram.py
Generates a clean, simple, fully-editable Use Case Diagram for the
V5 Hybrid Fraud Detection System.

Outputs:
  - PNG image  →  Project_Diagrams_PNG/Use_Case_Diagram.png
  - PPTX file  →  Project_Diagrams_PNG/Use_Case_Diagram.pptx  (fully editable native shapes)

Color palette and styling matches the existing Project_Diagrams.pptx theme.
"""

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = Path(__file__).resolve().parents[1]
PNG_DIR  = BASE_DIR / "Project_Diagrams_PNG"
PPTX_OUT = PNG_DIR / "Use_Case_Diagram.pptx"
PNG_OUT  = PNG_DIR / "Use_Case_Diagram.png"
PNG_DIR.mkdir(exist_ok=True)

# ── Colour palette (matches Project_Diagrams.pptx) ──────────────────────
C = {
    "bg":           "#FAFBFC",
    "primary":      "#2C3E50",
    "secondary":    "#34495E",
    "accent1":      "#2980B9",
    "accent2":      "#27AE60",
    "accent3":      "#E74C3C",
    "accent4":      "#F39C12",
    "accent5":      "#8E44AD",
    "accent6":      "#1ABC9C",
    "light_blue":   "#D6EAF8",
    "light_green":  "#D5F5E3",
    "light_red":    "#FADBD8",
    "light_orange": "#FDEBD0",
    "light_purple": "#E8DAEF",
    "light_teal":   "#D1F2EB",
    "light_gray":   "#F2F3F4",
    "border":       "#BDC3C7",
    "white":        "#FFFFFF",
    "line":         "#7F8C8D",
}

def _hex_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _hex_tuple(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2],16)/255.0 for i in (0,2,4))


# ═══════════════════════════════════════════════════════════════════════════
#  1.  MATPLOTLIB  (high-quality PNG)
# ═══════════════════════════════════════════════════════════════════════════

def _draw_actor_mpl(ax, cx, cy, label):
    """Draw a stick-figure actor with label below."""
    head_r = 2.5
    ax.add_patch(plt.Circle((cx, cy+12), head_r, fc=C["white"],
                            ec=C["primary"], lw=1.8, zorder=4))
    ax.plot([cx, cx],         [cy+9.5, cy+2],  color=C["primary"], lw=1.8, zorder=4)
    ax.plot([cx-5.5, cx+5.5], [cy+7, cy+7],    color=C["primary"], lw=1.8, zorder=4)
    ax.plot([cx, cx-4.5],     [cy+2, cy-5],     color=C["primary"], lw=1.8, zorder=4)
    ax.plot([cx, cx+4.5],     [cy+2, cy-5],     color=C["primary"], lw=1.8, zorder=4)
    ax.text(cx, cy-9, label, ha="center", va="top", fontsize=9,
            fontweight="bold", color=C["primary"], family="sans-serif")


def _draw_oval_mpl(ax, cx, cy, w, h, text, fill, edge):
    """Draw a use-case ellipse with centered multi-line text."""
    ell = plt.matplotlib.patches.Ellipse(
        (cx, cy), w, h, facecolor=fill, edgecolor=edge,
        linewidth=1.5, zorder=3)
    ax.add_patch(ell)
    ax.text(cx, cy, text, ha="center", va="center", fontsize=7.5,
            fontweight="bold", color=C["primary"], family="sans-serif",
            linespacing=1.35, zorder=5)


def _draw_arrow_mpl(ax, x1, y1, x2, y2, color, label="", dashed=False):
    ls = "--" if dashed else "-"
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle="->", color=color, lw=1.4,
                                linestyle=ls),
                zorder=2)
    if label:
        mx, my = (x1+x2)/2, (y1+y2)/2
        ax.text(mx, my+2.5, label, ha="center", va="center", fontsize=6,
                fontstyle="italic", color=C["secondary"], family="sans-serif",
                bbox=dict(boxstyle="round,pad=0.15", fc="white", ec="none",
                          alpha=0.9),
                zorder=6)


def _draw_line_mpl(ax, x1, y1, x2, y2, color):
    ax.plot([x1, x2], [y1, y2], color=color, lw=1.4, zorder=2)


def generate_png():
    """Render the Use Case diagram to a high-res PNG."""
    fig, ax = plt.subplots(figsize=(18, 12))
    fig.patch.set_facecolor(C["bg"])
    ax.set_xlim(-5, 185)
    ax.set_ylim(-15, 115)
    ax.set_aspect("equal")
    ax.axis("off")

    # Title
    ax.text(90, 110, "Use Case Diagram — V5 Hybrid Fraud Detection System",
            ha="center", va="center", fontsize=16, fontweight="bold",
            color=C["primary"], family="sans-serif")

    # ── System boundary (dashed rectangle) ──
    boundary = FancyBboxPatch((30, -5), 120, 100,
                              boxstyle="round,pad=2",
                              facecolor="#F8F9FA", edgecolor=C["primary"],
                              linewidth=1.8, linestyle="--", alpha=0.6, zorder=1)
    ax.add_patch(boundary)
    ax.text(90, 98, "V5 Hybrid Fraud Detection System",
            ha="center", va="center", fontsize=11, fontweight="bold",
            color=C["primary"], family="sans-serif",
            bbox=dict(boxstyle="round,pad=0.3", fc="white", ec=C["primary"],
                      lw=1), zorder=6)

    # ── Actors (outside boundary) ──
    _draw_actor_mpl(ax, 10, 65, "Transaction\nSource")
    _draw_actor_mpl(ax, 10, 15, "System\nAdmin")
    _draw_actor_mpl(ax, 170, 65, "Fraud\nAnalyst")

    # ── Use Cases (inside boundary) ──
    ow, oh = 30, 12  # oval width, height

    # Left column — Transaction Source use cases
    uc1_x, uc1_y = 58, 78   # Submit Transaction
    uc2_x, uc2_y = 58, 55   # Upload Batch CSV

    # Centre — System core use cases
    uc3_x, uc3_y = 90, 66   # Run V5 Fraud Detection
    uc4_x, uc4_y = 90, 35   # Generate Decision Report

    # Right column — Fraud Analyst use cases
    uc5_x, uc5_y = 122, 78  # Review Flagged Alerts
    uc6_x, uc6_y = 122, 55  # View Detection Results

    # Bottom — Admin use case
    uc7_x, uc7_y = 58, 12   # Configure Thresholds
    uc8_x, uc8_y = 122, 12  # Monitor Model Performance

    _draw_oval_mpl(ax, uc1_x, uc1_y, ow, oh, "Submit\nTransaction",   C["light_blue"],   C["accent1"])
    _draw_oval_mpl(ax, uc2_x, uc2_y, ow, oh, "Upload\nBatch CSV",     C["light_blue"],   C["accent1"])
    _draw_oval_mpl(ax, uc3_x, uc3_y, ow+4, oh+2, "Run V5 Fraud\nDetection", C["light_red"], C["accent3"])
    _draw_oval_mpl(ax, uc4_x, uc4_y, ow, oh, "Generate Decision\nReport", C["light_green"], C["accent2"])
    _draw_oval_mpl(ax, uc5_x, uc5_y, ow, oh, "Review Flagged\nAlerts", C["light_orange"], C["accent4"])
    _draw_oval_mpl(ax, uc6_x, uc6_y, ow, oh, "View Detection\nResults",C["light_teal"],   C["accent6"])
    _draw_oval_mpl(ax, uc7_x, uc7_y, ow, oh, "Configure\nThresholds", C["light_purple"],  C["accent5"])
    _draw_oval_mpl(ax, uc8_x, uc8_y, ow, oh, "Monitor Model\nPerformance", C["light_purple"], C["accent5"])

    # ── Actor → Use Case connections (solid lines) ──
    _draw_line_mpl(ax, 18, 70, uc1_x-15, uc1_y)    # Tx Source → Submit
    _draw_line_mpl(ax, 18, 60, uc2_x-15, uc2_y)    # Tx Source → Upload
    _draw_line_mpl(ax, 18, 20, uc7_x-15, uc7_y)    # Admin → Configure
    _draw_line_mpl(ax, 18, 15, uc8_x-15, uc8_y)    # Admin → Monitor
    _draw_line_mpl(ax, 162, 70, uc5_x+15, uc5_y)   # Analyst → Review
    _draw_line_mpl(ax, 162, 60, uc6_x+15, uc6_y)   # Analyst → View Results

    # ── <<include>> relationships (dashed arrows) ──
    _draw_arrow_mpl(ax, uc1_x+12, uc1_y-4, uc3_x-14, uc3_y+4,
                    C["accent5"], "<<include>>", dashed=True)
    _draw_arrow_mpl(ax, uc2_x+12, uc2_y+4, uc3_x-14, uc3_y-4,
                    C["accent5"], "<<include>>", dashed=True)

    # ── <<extend>> relationships (dashed arrows) ──
    _draw_arrow_mpl(ax, uc3_x+14, uc3_y+4, uc5_x-12, uc5_y-2,
                    C["accent4"], "<<extend>>", dashed=True)
    _draw_arrow_mpl(ax, uc3_x, uc3_y-7, uc4_x, uc4_y+7,
                    C["accent2"], "<<include>>", dashed=True)

    # Save
    fig.savefig(str(PNG_OUT), dpi=200, bbox_inches="tight", facecolor=C["bg"])
    plt.close(fig)
    print(f"PNG saved → {PNG_OUT}")


# ═══════════════════════════════════════════════════════════════════════════
#  2.  PPTX  (fully editable native shapes)
# ═══════════════════════════════════════════════════════════════════════════
SW, SH = 13.333, 7.5   # slide dimensions in inches

def _in(v):
    return Inches(v)

def _fmt_shape(shape, text, fill, text_color, edge_color, font_size,
               bold=True, italic=False):
    """Apply fill, border, and text formatting to a PPTX shape."""
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_rgb(fill)
    else:
        shape.fill.background()
    if edge_color and edge_color != "none":
        shape.line.color.rgb = _hex_rgb(edge_color)
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        try:
            from pptx.enum.text import MSO_ANCHOR
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.color.rgb = _hex_rgb(text_color)
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = PP_ALIGN.CENTER


def _add_oval(slide, cx, cy, w, h, text, fill, edge, fs=9):
    """Add an editable oval (use case) to the slide."""
    left = _in(cx - w/2)
    top  = _in(cy - h/2)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, _in(w), _in(h))
    _fmt_shape(shape, text, fill, C["primary"], edge, fs)
    return shape


def _add_box(slide, cx, cy, w, h, text, fill, edge, fs=9, bold=True):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _in(cx - w/2), _in(cy - h/2), _in(w), _in(h))
    _fmt_shape(shape, text, fill, C["primary"], edge, fs, bold=bold)
    return shape


def _add_connector(slide, x1, y1, x2, y2, color, arrow=False, dashed=False):
    """Add an editable connector line."""
    conn = slide.shapes.add_connector(
        1, _in(x1), _in(y1), _in(x2), _in(y2))
    conn.line.color.rgb = _hex_rgb(color)
    conn.line.width = Pt(1.5)
    if arrow:
        conn.line.end_arrowhead = 2
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return conn


def _add_label(slide, cx, cy, text, fs=7):
    """Add a small floating text label (e.g. <<include>>)."""
    w, h = 0.95, 0.28
    tb = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _in(cx - w/2), _in(cy - h/2), _in(w), _in(h))
    _fmt_shape(tb, text, C["white"], C["secondary"], "none", fs,
               bold=False, italic=True)
    return tb


def _add_actor_pptx(slide, cx, cy, label):
    """Draw a stick-figure actor from editable PPTX shapes."""
    # Head (small circle)
    head_d = 0.35
    slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        _in(cx - head_d/2), _in(cy - 0.17), _in(head_d), _in(head_d)
    ).line.color.rgb = _hex_rgb(C["primary"])

    # Body
    _add_connector(slide, cx, cy + head_d - 0.12, cx, cy + 0.65, C["primary"])
    # Arms
    _add_connector(slide, cx - 0.3, cy + 0.35, cx + 0.3, cy + 0.35, C["primary"])
    # Left leg
    _add_connector(slide, cx, cy + 0.65, cx - 0.25, cy + 1.0, C["primary"])
    # Right leg
    _add_connector(slide, cx, cy + 0.65, cx + 0.25, cy + 1.0, C["primary"])
    # Label
    tb = slide.shapes.add_textbox(
        _in(cx - 0.6), _in(cy + 1.05), _in(1.2), _in(0.45))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = _hex_rgb(C["primary"])
    p.alignment = PP_ALIGN.CENTER


def generate_pptx():
    """Build the fully-editable PPTX Use Case Diagram."""
    prs = Presentation()
    prs.slide_width  = _in(SW)
    prs.slide_height = _in(SH)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)

    # ── Title header bar ──
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _in(0), _in(0), _in(SW), _in(0.65))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = _hex_rgb(C["primary"])
    hdr.line.fill.background()
    p = hdr.text_frame.paragraphs[0]
    p.text = "Use Case Diagram — V5 Hybrid Fraud Detection System"
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ── Footer ──
    ft = slide.shapes.add_textbox(_in(0), _in(SH - 0.35), _in(SW), _in(0.3))
    fp = ft.text_frame.paragraphs[0]
    fp.text = "V5 Hybrid Fraud Detection System  ·  Major Project Documentation"
    fp.font.color.rgb = _hex_rgb(C["line"])
    fp.font.size = Pt(8)
    fp.alignment = PP_ALIGN.CENTER

    # ── System boundary (dashed rounded rectangle) ──
    bnd_x, bnd_y, bnd_w, bnd_h = 2.6, 0.85, 8.1, 6.0
    bnd = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _in(bnd_x), _in(bnd_y), _in(bnd_w), _in(bnd_h))
    bnd.fill.solid()
    bnd.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    bnd.fill.transparency = 0.4
    bnd.line.color.rgb = _hex_rgb(C["primary"])
    bnd.line.width = Pt(2)
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    bnd.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    # Boundary title
    _add_box(slide, bnd_x + bnd_w/2, bnd_y + 0.22, 3.6, 0.35,
             "V5 Hybrid Fraud Detection System",
             C["white"], C["primary"], fs=10)

    # ── Actors ──
    _add_actor_pptx(slide, 1.2, 2.2, "Transaction\nSource")
    _add_actor_pptx(slide, 1.2, 5.2, "System\nAdmin")
    _add_actor_pptx(slide, 12.1, 3.2, "Fraud\nAnalyst")

    # ── Use Cases (ovals) ──
    #    Carefully positioned to avoid any overlap.
    #    Layout: two columns + centre column

    ow, oh = 2.2, 0.95  # oval dimensions

    # Left column (Transaction Source use cases)
    uc1 = _add_oval(slide, 4.2, 2.1, ow, oh, "Submit\nTransaction",      C["light_blue"],  C["accent1"], 8)
    uc2 = _add_oval(slide, 4.2, 3.8, ow, oh, "Upload\nBatch CSV",        C["light_blue"],  C["accent1"], 8)

    # Centre column (Core system)
    uc3 = _add_oval(slide, 6.65, 2.95, ow+0.4, oh+0.15, "Run V5 Fraud\nDetection", C["light_red"], C["accent3"], 8)
    uc4 = _add_oval(slide, 6.65, 5.1, ow, oh, "Generate Decision\nReport", C["light_green"], C["accent2"], 8)

    # Right column (Fraud Analyst use cases)
    uc5 = _add_oval(slide, 9.1, 2.1, ow, oh, "Review Flagged\nAlerts",    C["light_orange"], C["accent4"], 8)
    uc6 = _add_oval(slide, 9.1, 3.8, ow, oh, "View Detection\nResults",   C["light_teal"],  C["accent6"], 8)

    # Bottom row (Admin use cases)
    uc7 = _add_oval(slide, 4.7, 5.8, ow, oh, "Configure\nThresholds",    C["light_purple"], C["accent5"], 8)
    uc8 = _add_oval(slide, 8.5, 5.8, ow, oh, "Monitor Model\nPerformance", C["light_purple"], C["accent5"], 8)

    # ── Actor → Use Case lines (solid, no arrowhead) ──
    _add_connector(slide, 1.8, 2.7,  4.2 - ow/2, 2.1,  C["line"])
    _add_connector(slide, 1.8, 3.0,  4.2 - ow/2, 3.8,  C["line"])
    _add_connector(slide, 1.8, 5.7,  4.7 - ow/2, 5.8,  C["line"])
    _add_connector(slide, 1.8, 5.9,  8.5 - ow/2, 5.8,  C["line"])
    _add_connector(slide, 11.5, 3.7, 9.1 + ow/2, 2.1,  C["line"])
    _add_connector(slide, 11.5, 4.0, 9.1 + ow/2, 3.8,  C["line"])

    # ── <<include>> dashed arrows ──
    # Submit Transaction --include--> Run V5
    _add_connector(slide, 4.2 + ow/2, 2.3, 6.65 - (ow+0.4)/2, 2.7,
                   C["accent5"], arrow=True, dashed=True)
    _add_label(slide, 5.25, 2.15, "<<include>>")

    # Upload CSV --include--> Run V5
    _add_connector(slide, 4.2 + ow/2, 3.6, 6.65 - (ow+0.4)/2, 3.15,
                   C["accent5"], arrow=True, dashed=True)
    _add_label(slide, 5.25, 3.65, "<<include>>")

    # Run V5 --include--> Generate Report
    _add_connector(slide, 6.65, 2.95 + (oh+0.15)/2, 6.65, 5.1 - oh/2,
                   C["accent2"], arrow=True, dashed=True)
    _add_label(slide, 6.65, 4.05, "<<include>>")

    # ── <<extend>> dashed arrows ──
    # Run V5 --extend--> Review Flagged Alerts
    _add_connector(slide, 6.65 + (ow+0.4)/2, 2.7, 9.1 - ow/2, 2.3,
                   C["accent4"], arrow=True, dashed=True)
    _add_label(slide, 8.05, 2.15, "<<extend>>")

    # Save
    prs.save(str(PPTX_OUT))
    print(f"PPTX saved → {PPTX_OUT}")


# ═══════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("="*55)
    print("  GENERATING USE CASE DIAGRAM  (PNG + Editable PPTX)")
    print("="*55)
    generate_png()
    generate_pptx()
    print("Done!")
